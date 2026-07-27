"""
File Type Router
================

Centralized file type classification and routing for the RAG pipeline.
Determines the appropriate processing method for each document type.
"""

import threading
from dataclasses import dataclass
from enum import Enum
from pathlib import Path
from typing import List

import logging
logger = logging.getLogger("FileTypeRouter")




class DocumentType(Enum):
    """Document type classification."""

    PDF = "pdf"
    TEXT = "text"
    MARKDOWN = "markdown"
    DOCX = "docx"
    PPTX = "pptx"
    IMAGE = "image"
    UNKNOWN = "unknown"


@dataclass
class FileClassification:
    """Result of file classification."""

    parser_files: List[str]
    text_files: List[str]
    docx_files: List[str]
    pptx_files: List[str]
    image_files: List[str]
    unsupported: List[str]


def serialize_table(tbl) -> str:
    """序列化 python-docx / python-pptx 表格为 Markdown 表格纯文本。

    对向量检索与 LLM 理解均友好。正确处理合并单元格：python-docx 的 row.cells
    对横向/纵向合并格返回同一底层 <w:tc> 的重复引用，按 id(cell._tc) 去重——
    横向合并的值只在首列出现一次，纵向合并的延续行该列留空，从而保持网格列对齐
    （治旧实现把合并表序列化成「标称值 / 实测值 | 750 | 580…」错乱的病）。pptx
    无底层 _tc，退用 id(cell) 尽力去重。docx/pptx 共用此实现。
    """
    rows = tbl.rows
    if not rows:
        return ""
    is_docx = all(hasattr(c, "_tc") for c in rows[0].cells)
    total_cols = max((len(r.cells) for r in rows), default=0)
    if total_cols == 0:
        return ""

    grid_rows: list[dict[int, str]] = []
    for row in rows:
        seen: set[int] = set()
        row_cells: dict[int, str] = {}
        for col_idx, cell in enumerate(row.cells):
            tc = cell._tc if is_docx else None
            key = id(tc) if is_docx else id(cell)
            if key in seen:
                continue
            seen.add(key)
            # 纵向合并延续格（vMerge=continue）：值在 restart 格，此处留空保对齐
            if is_docx and tc.vMerge == "continue":
                row_cells[col_idx] = ""
            else:
                # 单元格内换行/多空白会破坏 Markdown 表格行结构，规范为单空格
                row_cells[col_idx] = " ".join(cell.text.split())
        if any(row_cells.values()):
            grid_rows.append(row_cells)
    if not grid_rows:
        return ""

    lines: list[str] = []
    for r_idx, row_cells in enumerate(grid_rows):
        cells_out = [row_cells.get(c, "") for c in range(total_cols)]
        lines.append("| " + " | ".join(cells_out) + " |")
        if r_idx == 0:
            lines.append("| " + " | ".join(["---"] * total_cols) + " |")
    return "\n".join(lines)


# ── Docling 单例（PDF 解析 backend，arq worker 并发下只加载一次）──────────────
_DOCLING_CONVERTER = None
_DOCLING_CONVERTER_LOCK = threading.Lock()


def _get_docling_converter():
    """模块级 DocumentConverter 单例（首次加载用 Lock 守护）。

    读 settings.pdf（do_ocr / ocr_provider）配置 PdfPipelineOptions。docling 未装时
    抛 ImportError，由调用方 _extract_pdf_sections_docling 捕获 → 返回 []（不降级）。
    配置变更需重启进程（单例只建一次）。
    """
    global _DOCLING_CONVERTER
    if _DOCLING_CONVERTER is None:
        with _DOCLING_CONVERTER_LOCK:
            if _DOCLING_CONVERTER is None:
                from settings import get_settings

                cfg = get_settings().pdf
                from docling.document_converter import (
                    DocumentConverter,
                    PdfFormatOption,
                )
                from docling.datamodel.pipeline_options import (
                    PdfPipelineOptions,
                    RapidOcrOptions,
                )

                pipeline = PdfPipelineOptions()
                pipeline.do_ocr = bool(cfg.do_ocr)
                if cfg.ocr_provider == "rapid":
                    pipeline.ocr_options = RapidOcrOptions()
                _DOCLING_CONVERTER = DocumentConverter(
                    format_options={"pdf": PdfFormatOption(pipeline_options=pipeline)}
                )
    return _DOCLING_CONVERTER


def _docling_item_page(item) -> int:
    """从 Docling item.prov 取首个页码（1-based）；无 provenance 返回 0。"""
    prov = getattr(item, "prov", None) or []
    for p in prov:
        page_no = getattr(p, "page_no", None)
        if page_no:
            return int(page_no)
    return 0


class FileTypeRouter:
    """File type router for the RAG pipeline.

    Classifies files before processing to route them to appropriate handlers:
    - PDF files -> PDF parsing
    - Text files -> Direct read (fast, simple)
    - Unsupported -> Skip with warning
    """

    PARSER_EXTENSIONS = {".pdf"}

    TEXT_EXTENSIONS = {
        ".txt",
        ".text",
        ".log",
        ".md",
        ".markdown",
        ".rst",
        ".asciidoc",
        ".json",
        ".yaml",
        ".yml",
        ".toml",
        ".csv",
        ".tsv",
        ".tex",
        ".latex",
        ".bib",
        ".py",
        ".js",
        ".ts",
        ".jsx",
        ".tsx",
        ".java",
        ".c",
        ".cpp",
        ".h",
        ".hpp",
        ".go",
        ".rs",
        ".rb",
        ".php",
        ".swift",
        ".kt",
        ".scala",
        ".r",
        ".sql",
        ".sh",
        ".bash",
        ".zsh",
        ".ps1",
        ".html",
        ".htm",
        ".xml",
        ".css",
        ".scss",
        ".sass",
        ".less",
        ".ini",
        ".cfg",
        ".conf",
        ".env",
        ".properties",
    }

    DOCX_EXTENSIONS = {".docx"}
    PPTX_EXTENSIONS = {".pptx"}
    IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif"}

    # M-27：legacy 二进制格式（.doc/.ppt）仅支持 OOXML 解析（python-docx/python-pptx），
    # 这两类无法解析。上传层应已拒绝（见 api/admin._ALLOWED_EXT）；此处列出用于在
    # classify_files 的 unsupported 分支给出针对性提示，便于排查上游漏网情况。
    LEGACY_UNSUPPORTED = {".doc", ".ppt", ".xls", ".wps"}

    @classmethod
    def get_document_type(cls, file_path: str) -> DocumentType:
        """Classify a single file by its type."""
        ext = Path(file_path).suffix.lower()

        if ext in cls.PARSER_EXTENSIONS:
            return DocumentType.PDF
        elif ext in cls.TEXT_EXTENSIONS:
            return DocumentType.TEXT
        elif ext in cls.DOCX_EXTENSIONS:
            return DocumentType.DOCX
        elif ext in cls.PPTX_EXTENSIONS:
            return DocumentType.PPTX
        elif ext in cls.IMAGE_EXTENSIONS:
            return DocumentType.IMAGE
        else:
            if cls._is_text_file(file_path):
                return DocumentType.TEXT
            return DocumentType.UNKNOWN

    @classmethod
    def _is_text_file(cls, file_path: str, sample_size: int = 8192) -> bool:
        """Detect if a file is text-based by examining its content."""
        try:
            with open(file_path, "rb") as f:
                chunk = f.read(sample_size)

            if b"\x00" in chunk:
                return False

            chunk.decode("utf-8")
            return True
        except (UnicodeDecodeError, IOError, OSError):
            return False

    @classmethod
    def classify_files(cls, file_paths: List[str]) -> FileClassification:
        """Classify a list of files by processing method."""
        parser_files = []
        text_files = []
        docx_files = []
        pptx_files = []
        image_files = []
        unsupported = []

        for path in file_paths:
            doc_type = cls.get_document_type(path)

            if doc_type == DocumentType.PDF:
                parser_files.append(path)
            elif doc_type in (DocumentType.TEXT, DocumentType.MARKDOWN):
                text_files.append(path)
            elif doc_type == DocumentType.DOCX:
                docx_files.append(path)
            elif doc_type == DocumentType.PPTX:
                pptx_files.append(path)
            elif doc_type == DocumentType.IMAGE:
                image_files.append(path)
            else:
                unsupported.append(path)

        logger.debug(
            f"Classified {len(file_paths)} files: "
            f"{len(parser_files)} parser, {len(text_files)} text, {len(docx_files)} docx, "
            f"{len(pptx_files)} pptx, {len(image_files)} image, {len(unsupported)} unsupported"
        )
        # M-27：识别 legacy 格式（.doc/.ppt 等），给针对性 warning——这类文件无解析
        # handler，会被静默丢弃。若此处出现，说明上游上传校验漏网。
        legacy_hits = [p for p in unsupported if Path(p).suffix.lower() in cls.LEGACY_UNSUPPORTED]
        if legacy_hits:
            logger.warning(
                "检测到 legacy 格式文件（无解析 handler，将被跳过）：%s。"
                "请转换为 .docx/.pptx 后重新上传。",
                [Path(p).name for p in legacy_hits],
            )

        return FileClassification(
            parser_files=parser_files,
            text_files=text_files,
            docx_files=docx_files,
            pptx_files=pptx_files,
            image_files=image_files,
            unsupported=unsupported,
        )

    @classmethod
    def extract_docx_text(cls, file_path: str) -> str:
        """Read plain text from a .docx (Office Open XML). Legacy .doc is not supported."""
        p = Path(file_path)
        if p.suffix.lower() != ".docx":
            return ""
        try:
            from docx import Document
        except ImportError:
            logger.warning("python-docx is not installed; cannot read .docx")
            return ""
        try:
            d = Document(str(p))
            parts: list[str] = []
            for para in d.paragraphs:
                t = (para.text or "").strip()
                if t:
                    parts.append(t)
            for table in d.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            return "\n\n".join(parts)
        except Exception as exc:
            logger.warning("Failed to read .docx %s: %s", p.name, exc)
            return ""

    @classmethod
    def extract_docx_sections(cls, file_path: str) -> list[dict]:
        """Parse a .docx and split it into sections by Heading 1.

        Each returned dict has:
          - title   : heading text (e.g. "实验九 二端口网络研究")
          - content : full text of that section, including table cells with
                      column-header prefixes and [图: 电路图] image placeholders
          - metadata: {"section": title, "file_name": basename}

        Falls back to a single section containing the whole document when
        no Heading 1 paragraphs are found (e.g. plain text docs).
        """
        p = Path(file_path)
        if p.suffix.lower() != ".docx":
            return []
        try:
            from docx import Document as _DocxDocument
        except ImportError:
            logger.warning("python-docx is not installed; cannot read .docx")
            return []

        NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        DML = "http://schemas.openxmlformats.org/drawingml/2006/main"
        PIC = "http://schemas.openxmlformats.org/drawingml/2006/picture"

        try:
            d = _DocxDocument(str(p))
        except Exception as exc:
            logger.warning("Failed to open .docx %s: %s", p.name, exc)
            return []

        # 表格序列化用模块级 serialize_table（docx/pptx 共用，见文件顶部）

        # Walk body XML children to preserve paragraph/table order per section
        body = d.element.body
        sections: list[dict] = []
        current_title: str | None = None
        current_parts: list[str] = []

        # Map xml element id → table object for quick lookup
        tbl_map: dict[int, object] = {id(t._element): t for t in d.tables}

        def _flush(title: str | None, parts: list[str]) -> None:
            if parts:
                content = "\n\n".join(parts)
                t = title or ""
                sections.append({
                    "title": t,
                    "content": (f"{t}\n\n{content}").strip() if t else content,
                    "metadata": {"section": t, "file_name": p.name},
                })

        # pStyle val values that map to Heading 1 across different Word versions/locales:
        # "Heading1"  — English style id
        # "1"         — common numeric alias
        # "2"         — OOXML built-in numeric id observed in Chinese Word installs
        _H1_SVALS = {"Heading1", "1", "2"}

        for child in list(body):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

            if tag == "p":
                pStyle = child.find(f".//{{{NS}}}pStyle")
                sval = (pStyle.get(f"{{{NS}}}val") or "") if pStyle is not None else ""
                txt = "".join(t.text or "" for t in child.iter(f"{{{NS}}}t")).strip()

                # Heading 1 → start new section (require non-empty title)
                if sval in _H1_SVALS and txt:
                    _flush(current_title, current_parts)
                    current_title = txt
                    current_parts = []
                    continue

                # Image detection
                has_image = (
                    child.find(f".//{{{DML}}}blip") is not None
                    or child.find(f".//{{{PIC}}}pic") is not None
                )
                if has_image:
                    label = f"[图: {txt}]" if txt else "[图: 电路图]"
                    current_parts.append(label)
                elif txt:
                    current_parts.append(txt)

            elif tag == "tbl":
                tbl_obj = tbl_map.get(id(child))
                if tbl_obj is not None:
                    serialized = serialize_table(tbl_obj)
                    if serialized:
                        current_parts.append(serialized)

        _flush(current_title, current_parts)

        # If no Heading 1 was found, fall back to whole-doc single section
        if not sections:
            fallback = cls.extract_docx_text(file_path)
            if fallback:
                sections.append({
                    "title": p.stem,
                    "content": fallback,
                    "metadata": {"section": p.stem, "file_name": p.name},
                })

        logger.info(
            "extract_docx_sections: %s → %d sections", p.name, len(sections)
        )
        return sections

    @classmethod
    def _serialize_pptx_table(cls, table) -> str:
        """Serialize a python-pptx table（委托模块级 serialize_table，与 DOCX 同逻辑）。"""
        return serialize_table(table)

    @classmethod
    def extract_pptx_sections(cls, file_path: str) -> list[dict]:
        """Parse a .pptx and split it into one section per slide.

        Each returned dict has:
          - title   : slide title (or "Slide N")
          - content : title + body text and tables for that slide
          - metadata: {"section": title, "file_name": basename}
        """
        p = Path(file_path)
        if p.suffix.lower() != ".pptx":
            return []
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx is not installed; cannot read .pptx")
            return []

        try:
            prs = Presentation(str(p))
        except Exception as exc:
            logger.warning("Failed to open .pptx %s: %s", p.name, exc)
            return []

        sections: list[dict] = []
        for i, slide in enumerate(prs.slides):
            title_shape = slide.shapes.title
            title = ""
            if title_shape is not None and title_shape.has_text_frame:
                title = (title_shape.text or "").strip()
            if not title:
                title = f"Slide {i + 1}"

            body_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_table:
                    serialized = cls._serialize_pptx_table(shape.table)
                    if serialized:
                        body_parts.append(serialized)
                    continue
                if not shape.has_text_frame:
                    continue
                if title_shape is not None and shape is title_shape:
                    continue
                for para in shape.text_frame.paragraphs:
                    t = (para.text or "").strip()
                    if t and t != title:
                        body_parts.append(t)

            if not body_parts:
                continue

            content = "\n\n".join(body_parts)
            sections.append({
                "title": title,
                "content": (f"{title}\n\n{content}").strip(),
                "metadata": {"section": title, "file_name": p.name},
            })

        logger.info(
            "extract_pptx_sections: %s → %d sections", p.name, len(sections)
        )
        return sections

    @classmethod
    def read_text_file_sync(cls, file_path: str) -> str:
        """Read a text file with automatic encoding detection (sync; 与 LightRAG 线程池解析共用)."""
        encodings = ["utf-8", "utf-8-sig", "gbk", "gb2312", "gb18030", "latin-1", "cp1252"]

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue

        with open(file_path, "rb") as f:
            return f.read().decode("utf-8", errors="replace")

    @classmethod
    def extract_pdf_text(cls, file_path: str) -> str:
        """Extract PDF text with PyMuPDF（LlamaIndex / LightRAG 摄入共用）。"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.warning("PyMuPDF not installed. Cannot extract PDF text.")
            return ""
        try:
            doc = fitz.open(file_path)
            texts = []
            for page in doc:
                texts.append(page.get_text())
            doc.close()
            return "\n\n".join(texts)
        except Exception as exc:
            logger.warning("Failed to extract PDF text %s: %s", file_path, exc)
            return ""

    @classmethod
    def extract_pdf_sections(cls, file_path: str, backend: str = "docling") -> list[dict]:
        """Parse a .pdf → sections（与 extract_docx_sections 同构：[{title, content, page}]）。

        backend="docling"（默认）：Docling 单引擎全包（版面/表格/标题/页码 + 扫描件 OCR），
            表格原子化为独立 section（不锯断）；标题(label=section_header)开新 section。
        backend="mupdf"：PyMuPDF get_toc() 切章节 + page.get_text() + 页码注入，
            轻量纯 CPU（无 torch）；无 OCR、无表格原子化。
        选定 backend 失败 → logger.warning + 返回 []（该文件跳过，**不切到另一个 backend**）。
        要换 backend 改 settings.pdf.backend 重启，无运行时降级。
        """
        if backend == "mupdf":
            return cls._extract_pdf_sections_mupdf(file_path)
        return cls._extract_pdf_sections_docling(file_path)

    @classmethod
    def _extract_pdf_sections_mupdf(cls, file_path: str) -> list[dict]:
        """mupdf backend：fitz.get_toc() 切章节 + page.get_text() + 页码 → sections。"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.warning("PyMuPDF 未安装，PDF(mupdf backend) 跳过: %s", file_path)
            return []
        try:
            doc = fitz.open(file_path)
        except Exception as exc:
            logger.warning("mupdf 打开 PDF 失败 %s: %s", file_path, exc)
            return []
        try:
            toc = doc.get_toc()  # [[level, title, page_1based], ...]；无书签 → []
            if not toc:
                text = "\n\n".join(p.get_text() for p in doc)
                return [{"title": "", "content": text, "page": 1}] if text.strip() else []
            sections: list[dict] = []
            n = doc.page_count
            for i, (_lvl, title, start_page) in enumerate(toc):
                end_page = toc[i + 1][2] - 1 if i + 1 < len(toc) else n
                end_page = min(end_page, n)
                content = "\n\n".join(
                    doc[p].get_text() for p in range(start_page - 1, end_page)
                )
                if content.strip():
                    sections.append({"title": title, "content": content, "page": start_page})
            return sections
        except Exception as exc:
            logger.warning("mupdf 解析 PDF 失败 %s: %s", file_path, exc)
            return []
        finally:
            doc.close()

    @classmethod
    def _extract_pdf_sections_docling(cls, file_path: str) -> list[dict]:
        """docling backend：DocumentConverter → DoclingDocument → sections（表格原子化）。"""
        try:
            converter = _get_docling_converter()
        except ImportError as exc:
            logger.error(
                "backend=docling 但 docling 依赖未装，该 PDF 跳过（不降级 mupdf）%s: %s",
                file_path, exc,
            )
            return []
        try:
            result = converter.convert(file_path)
            ddoc = result.document
        except Exception as exc:
            logger.warning("docling 解析 PDF 失败 %s: %s", file_path, exc)
            return []

        sections: list[dict] = []
        current_title: str | None = None
        current_parts: list[str] = []
        current_page: int = 0

        def _flush() -> None:
            nonlocal current_parts, current_page
            if current_parts:
                content = "\n\n".join(current_parts)
                current_parts = []
                if content.strip():
                    sections.append({
                        "title": current_title or "",
                        "content": content,
                        "page": current_page,
                    })
            current_page = 0

        for item, _level in ddoc.iterate_items():
            label = getattr(item, "label", "")
            page = _docling_item_page(item) or current_page
            if label == "section_header":
                text = (getattr(item, "text", "") or "").strip()
                if text:
                    _flush()
                    current_title = text
                    current_page = page or current_page
            elif label == "table":
                _flush()
                try:
                    tbl_md = item.export_to_markdown(doc=ddoc)
                except Exception:
                    tbl_md = getattr(item, "text", "") or ""
                if tbl_md and tbl_md.strip():
                    sections.append({
                        "title": f"{current_title or ''}（表格）",
                        "content": tbl_md,
                        "page": page,
                    })
            else:  # paragraph / list_item / caption / ...
                text = (getattr(item, "text", "") or "").strip()
                if text:
                    current_parts.append(text)
                    if not current_page:
                        current_page = page
        _flush()

        if not sections:  # 无标题无表格 → 整文档单 section
            md = ddoc.export_to_markdown()
            if md and md.strip():
                sections.append({"title": "", "content": md, "page": 0})
        return sections

    @classmethod
    async def read_text_file(cls, file_path: str) -> str:
        """Read a text file with automatic encoding detection."""
        return cls.read_text_file_sync(file_path)

    @classmethod
    def needs_parser(cls, file_path: str) -> bool:
        """Quick check if a single file needs parser processing."""
        doc_type = cls.get_document_type(file_path)
        return doc_type in (
            DocumentType.PDF,
            DocumentType.DOCX,
            DocumentType.PPTX,
            DocumentType.IMAGE,
        )

    @classmethod
    def is_text_readable(cls, file_path: str) -> bool:
        """Check if a file can be read directly as text."""
        doc_type = cls.get_document_type(file_path)
        return doc_type in (DocumentType.TEXT, DocumentType.MARKDOWN)

    @classmethod
    def get_supported_extensions(cls) -> set[str]:
        """Get the set of all supported file extensions."""
        return (
            cls.PARSER_EXTENSIONS
            | cls.TEXT_EXTENSIONS
            | cls.DOCX_EXTENSIONS
            | cls.PPTX_EXTENSIONS
        )

    @classmethod
    def get_glob_patterns(cls) -> list[str]:
        """Get glob patterns for file searching."""
        return [f"*{ext}" for ext in sorted(cls.get_supported_extensions())]
